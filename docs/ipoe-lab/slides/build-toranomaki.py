#!/usr/bin/env python3
"""IPv6/IPoE 勉強会「虎の巻」PowerPoint 生成スクリプト

  usage: pip install python-pptx && python3 build-toranomaki.py [出力パス]

pip が使えない / システムに手を入れたくない環境 (Proxmox ホスト等) では、
wheel を落として展開し PYTHONPATH で読ませれば、パッケージを 1 つも
インストールせずに実行できる:

  mkdir -p ~/pptxlib ~/wheels && cd ~/wheels
  # python-pptx / lxml / pillow / XlsxWriter / typing_extensions の wheel を取得
  # (lxml と pillow は cp3XX-manylinux-x86_64 の版を選ぶこと)
  python3 -c "import zipfile,glob,os; [zipfile.ZipFile(w).extractall(os.path.expanduser('~/pptxlib')) for w in glob.glob('*.whl')]"
  PYTHONPATH=~/pptxlib python3 build-toranomaki.py

Proxmox ホストで apt から入れようとすると libc6 と dpkg の更新を
巻き込む (65 パッケージ)。ハイパーバイザでやってはいけない。

内容の正本は docs/ipoe-lab/study-guide.md。本スクリプトはそれを
プレゼン形式に落としたもの。内容を直す場合は study-guide.md も併せて直すこと。

**話す人向けの指示はスライド本体に書かないこと。**「〜と宣言してください」「生徒役は若手に」
のような講師への指示は **notes= に渡してノートへ**。投影すると聴衆に見えてしまいます。

フォントは Windows 標準の「Meiryo UI」を指定している (会社 PC で開く前提)。
Mac/Linux で開く場合は自動で代替フォントになる。
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

FONT = "Meiryo UI"
NAVY = RGBColor(0x1F, 0x33, 0x55)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
WARN = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "toranomaki.pptx")


def _run(paragraph, text):
    """段落に run を追加してテキストを入れる (add_run() は引数を取らない)"""
    r = paragraph.add_run()
    r.text = text
    return r


def _style(run, size, bold=False, color=NAVY):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _put(paragraph, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    if align is not None:
        paragraph.alignment = align
    r = _run(paragraph, text)
    _style(r, size, bold, color)
    return r


def _textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True   # 未指定だと箱が自動縮小され中央寄せに見える
    return tb


def _notes(s, notes):
    """**話す人向けの指示は必ずここへ。**
    スライド本体に「〜と宣言してください」「ここは 5 分かけて」と書くと、
    **投影されて聴衆に見えてしまう。**司会への指示はノートに置くこと。
    """
    if not notes:
        return s
    tf = s.notes_slide.notes_text_frame
    lines = notes if isinstance(notes, (list, tuple)) else [notes]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
    return s


def title_slide(prs, title, subtitle, lines, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = _textbox(s, 0.6, 1.6, 12.1, 1.4)
    tf = band.text_frame
    _put(tf.paragraphs[0], title, 40, True, NAVY)
    p = tf.add_paragraph()
    _put(p, subtitle, 18, False, ACCENT)

    box = s.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(12.0), Inches(3.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _put(p, ln, 14, False, GRAY)
        p.space_after = Pt(8)
    return _notes(s, notes)


def section_slide(prs, num, title, goal, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    box = _textbox(s, 1.0, 2.4, 11.3, 2.6)
    tf = box.text_frame
    tf.word_wrap = True
    _put(tf.paragraphs[0], num, 20, True, RGBColor(0x7F, 0xB8, 0xE0))
    p = tf.add_paragraph()
    _put(p, title, 36, True, WHITE)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    _put(p, "ゴール: " + goal, 16, False, RGBColor(0xD5, 0xE5, 0xF2))
    return _notes(s, notes)


def _text_width(text):
    """全角を 1.0、半角を 0.5 として文字幅を数える (日本語混在の折り返し見積もり用)"""
    w = 0.0
    for ch in text:
        w += 0.5 if ord(ch) < 0x2E80 else 1.0
    return w


def _wrapped_lines(text, size_pt, box_in):
    """そのフォントサイズ・箱幅で何行に折り返されるかの概算"""
    per_line = max(6.0, (box_in * 72.0) / (size_pt * 1.03))
    w = _text_width(text)
    return max(1, int(w / per_line) + (1 if w % per_line else 0))


def bullet_slide(prs, title, bullets, footer=None, notes=None):
    """bullets: [(indent, text, style)] style: '' / 'b'(太字) / 'w'(警告赤) / 'c'(コード)

    **箱からはみ出さないよう、行数を見積もってフォントと行間を自動で詰める。**
    枚数が増えても手で調整しなくて済むようにするため。既定の行間は広めに取ってある
    (読み上げながら目で追う資料なので、詰まっていると行を見失う)。
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _textbox(s, 0.55, 0.35, 12.2, 0.9)
    _put(tb.text_frame.paragraphs[0], title, 28, True, NAVY)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.22), Inches(12.2), Pt(2.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()

    BOX_H = 5.15 if footer else 5.45          # インチ。footer があるぶん本文を短くする
    BOX_W = 12.0

    def _plan(shrink):
        """このフォント縮小量で、本文が何インチになるかを返す"""
        total = 0.0
        for ind, text, _st in bullets:
            size = max(12, 20 - 3 * ind - shrink)
            gap = max(3, (10 if ind == 0 else 6) - shrink)
            lines = _wrapped_lines(text, size, BOX_W - 0.35 * ind)
            total += (lines * size * 1.30 + gap) / 72.0
        return total

    shrink = 0
    while shrink < 9 and _plan(shrink) > BOX_H:
        shrink += 1

    box = _textbox(s, 0.65, 1.5, BOX_W, BOX_H)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (ind, text, st) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = ind
        size = max(12, 20 - 3 * ind - shrink)
        color, bold = NAVY if ind == 0 else GRAY, False
        if st == "b":
            bold = True
        elif st == "w":
            color, bold = WARN, True
        elif st == "c":
            color = RGBColor(0x0B, 0x60, 0x40)
        r = _put(p, text, size, bold, color)
        if st == "c":
            r.font.name = "Consolas"
        p.space_after = Pt(max(3, (10 if ind == 0 else 6) - shrink))

    if footer:
        fb = _textbox(s, 0.65, 6.80, 12.0, 0.55)
        ftf = fb.text_frame
        ftf.word_wrap = True
        _put(ftf.paragraphs[0], footer, 11 if _text_width(footer) > 55 else 12, False, ACCENT)
    return _notes(s, notes)


def qa_slide(prs, wall_no, question, answer, dissent=None, ref=None, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _textbox(s, 0.55, 0.3, 12.2, 0.7)
    _put(tb.text_frame.paragraphs[0], f"つまずきの壁 {wall_no}", 22, True, ACCENT)

    y = 1.05
    # 生徒
    qb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y), Inches(12.2), Inches(1.05))
    qb.fill.solid()
    qb.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xFA)
    qb.line.color.rgb = ACCENT
    tf = qb.text_frame
    tf.word_wrap = True
    tf.margin_left, tf.margin_top = Inches(0.18), Inches(0.1)
    p = tf.paragraphs[0]
    _put(p, "生徒: ", 16, True, ACCENT)
    _put(p, question, 16, False, NAVY)

    y += 1.3
    ab = _textbox(s, 0.7, y, 12.0, 3.0)
    tf = ab.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(answer):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            _put(p, "先生: ", 16, True, NAVY)
            _put(p, ln, 16, False, NAVY)
        else:
            p.level = 1
            _put(p, ln, 14, False, GRAY)
        p.space_after = Pt(6)

    if dissent:
        db = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.15))
        db.fill.solid()
        db.fill.fore_color.rgb = RGBColor(0xFD, 0xF1, 0xEF)
        db.line.color.rgb = WARN
        tf = db.text_frame
        tf.word_wrap = True
        tf.margin_left, tf.margin_top = Inches(0.18), Inches(0.08)
        p = tf.paragraphs[0]
        _put(p, "対立意見(懐疑的な同僚): ", 13, True, WARN)
        _put(p, dissent, 13, False, GRAY)

    if ref:
        fb = _textbox(s, 0.7, 6.45, 11.9, 0.6)
        tf = fb.text_frame
        tf.word_wrap = True
        _put(tf.paragraphs[0], "教科書 / ラボ: " + ref, 12, False, ACCENT)
    return _notes(s, notes)


def table_slide(prs, title, headers, rows, widths, footer=None, fsize=12, notes=None):
    """セルの値が "!" で始まると、その 1 セルだけ赤の太字で強調する ("!" は表示されない)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _textbox(s, 0.55, 0.3, 12.2, 0.8)
    _put(tb.text_frame.paragraphs[0], title, 26, True, NAVY)

    nrow, ncol = len(rows) + 1, len(headers)

    # **セルが折り返すと PowerPoint が行を勝手に伸ばし、表がスライドからはみ出す。**
    # 行数と、各セルが列幅に対して何行になるかを見積もって、収まるまで文字を縮める。
    AVAIL = 5.5   # インチ。表に使える高さ
    while fsize > 8:
        est = 0.0
        for r, row in enumerate([headers] + rows):
            ln = 1
            for c, val in enumerate(row):
                w = widths[c] if c < len(widths) else 2.0
                ln = max(ln, _wrapped_lines(str(val).lstrip("!"), fsize, max(0.8, w - 0.25)))
            est += max(0.34, ln * fsize * 1.45 / 72.0)
        if est <= AVAIL:
            break
        fsize -= 1

    h = min(AVAIL, max(0.34 * nrow, est if 'est' in dir() else 0.42 * nrow))
    gf = s.shapes.add_table(nrow, ncol, Inches(0.55), Inches(1.25), Inches(12.2), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.last_row = False
    tbl.first_col = False
    tbl.last_col = False
    tbl.horz_banding = False
    tbl.vert_banding = False
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        _put(cell.text_frame.paragraphs[0], htxt, fsize + 1, True, WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            warn = val.startswith("!")
            _put(p, val.lstrip("!"), fsize, warn, WARN if warn else NAVY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF6, 0xFA)
    if footer:
        # 表が下に伸びても重ならないよう、見積もった高さの下に余白を空けて置く。
        # est は上のループで求めた「折り返しを含む実際の高さ」なので、これを使う。
        top = min(6.75, 1.25 + max(h, est if "est" in dir() else h) + 0.28)
        fb = _textbox(s, 0.55, top, 12.2, 0.6)
        tf = fb.text_frame
        tf.word_wrap = True
        _put(tf.paragraphs[0], footer, 12, False, ACCENT)
    return _notes(s, notes)


def branch_slide(prs, title, root, left, right, note=None, notes=None):
    """root から left / right へ分岐する図。root/left/right は (label, sublabel)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _textbox(s, 0.55, 0.3, 12.2, 0.8)
    _put(tb.text_frame.paragraphs[0], title, 26, True, NAVY)

    def _box(x, y, w, h, label, sub, color):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xFA)
        sh.line.color.rgb = color
        sh.line.width = Pt(1.75)
        tf = sh.text_frame
        tf.word_wrap = True
        _put(tf.paragraphs[0], label, 15, True, NAVY, PP_ALIGN.CENTER)
        _put(tf.add_paragraph(), sub, 11, False, GRAY, PP_ALIGN.CENTER)
        return sh

    _box(5.0, 1.3, 3.3, 1.0, root[0], root[1], NAVY)
    _box(1.3, 3.2, 4.6, 1.2, left[0], left[1], ACCENT)
    _box(7.4, 3.2, 4.6, 1.2, right[0], right[1], ACCENT)
    def _line(x1, y1, x2, y2):
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        cn.line.color.rgb = ACCENT
        cn.line.width = Pt(2.5)

    _line(6.65, 2.3, 6.65, 2.75)          # 分岐点から下へ
    _line(3.6, 2.75, 9.7, 2.75)           # 横に伸ばす
    _line(3.6, 2.75, 3.6, 3.2)            # 左の箱へ
    _line(9.7, 2.75, 9.7, 3.2)            # 右の箱へ

    if note:
        nb = _textbox(s, 0.75, 4.75, 11.8, 2.0)
        tf = nb.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(note):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _put(p, ln, 14, ln.startswith("※"), WARN if ln.startswith("※") else GRAY)
            p.space_after = Pt(6)
    return _notes(s, notes)


def diagram_slide(prs, title, boxes, note=None, notes=None):
    """boxes: [(label, sublabel)] を横一列に並べて矢印でつなぐ"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _textbox(s, 0.55, 0.3, 12.2, 0.8)
    _put(tb.text_frame.paragraphs[0], title, 26, True, NAVY)

    n = len(boxes)
    bw = 11.8 / n * 0.82
    gap = (11.8 - bw * n) / max(1, n - 1)
    x = 0.75
    for i, (label, sub) in enumerate(boxes):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.6), Inches(bw), Inches(1.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xEF, 0xF5, 0xFA)
        sh.line.color.rgb = ACCENT
        sh.line.width = Pt(1.75)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _put(p, label, 15, True, NAVY, PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        _put(p2, sub, 11, False, GRAY, PP_ALIGN.CENTER)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw), Inches(3.2), Inches(gap), Inches(0.32))
            ar.fill.solid()
            ar.fill.fore_color.rgb = ACCENT
            ar.line.fill.background()
        x += bw + gap

    if note:
        nb = _textbox(s, 0.75, 4.5, 11.8, 2.0)
        tf = nb.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(note):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _put(p, ln, 14, ln.startswith("※"), WARN if ln.startswith("※") else GRAY)
            p.space_after = Pt(6)
    return _notes(s, notes)

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # ---------------- 表紙 ----------------
    title_slide(
        prs,
        "IPv6 / IPoE 虎の巻",
        "PPPoE → IPoE 切替でトラブらないための 3 回シリーズ",
        [
            "対象: IPv6 をほとんど触ったことがない若手 (支店配属 1〜3 年目想定)",
            "教科書: 小川晃通『プロフェッショナル IPv6 第2版』(ラムダノート / CC BY-NC-SA)",
            "検証環境: IPoE 検証ラボ (構築中。自宅プロトタイプ検証後に社内展開)",
            "",
            "進め方: ラボで壊す → 「なんで?」が生まれる → 教科書の該当節を読む → もう一度ラボ",
            "",
            "社内限定資料。教科書の本文・図版は貼らず、該当ページを各自開いてください。",
        ],
    )

    bullet_slide(
        prs,
        "この勉強会の目的",
        [
            (0, "PPPoE → IPoE 切替案件が増えている。しかしトラブルも増えている", "b"),
            (1, "組み合わせが多い: IPv6 配布方式 × IPv4 over IPv6 方式 × CPE 機種 × 構成", ""),
            (1, "「やってみたら動かない」を客先で起こしている", "w"),
            (0, "原因は知識の欠落ではなく、原理を見たことがないこと", "b"),
            (1, "RA も DHCPv6-PD も MAP-E も、パケットを一度見れば腹に落ちる", ""),
            (1, "だから「検証環境 × 教科書」でやる。座学だけ / 手順書だけ、では身に付かない", ""),
            (0, "3 回終わったときにできるようになること", "b"),
            (1, "ひかり電話の有無と「どこで終端するか」の 2 段で方式を選べる", ""),
            (1, "「ポート開放できますか?」に方式を聞いて即答できる", ""),
            (1, "「Web は見えるがファイルが送れない」を切り分けて、直し方まで言える", ""),
        ],
        footer="正本: docs/ipoe-lab/study-guide.md",
    )

    bullet_slide(
        prs,
        "この教材の現在のステータス",
        [
            (0, "使える部分", "b"),
            (1, "教科書との対応表 / 用語集 / つまずきの壁 Q&A — ラボの状態に依存しない", ""),
            (0, "まだ確定していない部分", "w"),
            (1, "時間配分は仮案。実測していない。休憩・質疑の余白を入れていないので溢れる", ""),
            (1, "サイクル 1〜7 実走済み。R1〜R6/D1 の 7 件も実測完了", ""),
            (1, "実機 Cisco 892FJ での DS-Lite も実測済み (第4部を参照)", ""),
            (1, "会社 VMware 用ランブック (runbook-vmware.md) は執筆済み", ""),
            (0, "配布前にやること", "b"),
            (1, "proxmox-prototype.md §0.5 → §2 → §5 を一巡する", ""),
            (1, "実施前に test-matrix §4 の「⚠ 実測」注記を必ず読む (5 件が想定と違った)", ""),
            (1, "第1回はパイロット実施。実測が取れるまで第2回以降の日程を確定しない", ""),
        ],
        notes=[
            "【講師へ】このページは必読。**受講者には投影しなくてよい。**",
            "・時間配分は実測していない仮案。休憩と質疑の余白が無いので、そのままだと溢れる。",
            "・実施前に test-matrix §4 の「⚠ 実測」注記を必ず読むこと。5 件が想定と違った。",
            "・第1回はパイロットとして回し、実測が取れるまで第2回以降の日程を確定しないこと。",
        ],
    )

    table_slide(
        prs,
        "全体の流れ",
        ["回 / 部", "テーマ", "教科書", "ラボ演習"],
        [
            ["第0〜2部", "全体像 → IPv6 そのもの → 日本の回線 (座学・導入)", "1,2,3,5,7章", "—"],
            ["第1回", "なんで IPv6 に変えなきゃいけないの?", "1, 2, 18章", "演習 1-A"],
            ["第2回", "アドレスはどうやって降ってくるのか (最重要)", "6, 7, 8章, 付録B", "演習 2-A / 2-B"],
            ["第3回", "IPv4 を IPv6 の上で運ぶ", "24章", "演習 3-A (R3/R4/R5)"],
            ["補講", "セキュリティ / DNS / NGN 固有問題", "14, 17章, 付録A", "R6, R13"],
            ["第4部", "検証結果 (ラボ / 実機 892FJ / CML の実測値)", "—", "—"],
            ["巻末", "つまずきの壁 Q&A (各回 +15〜20 分)", "—", "—"],
        ],
        [1.1, 6.3, 2.4, 2.4],
        footer="全回で参照: 付録 A「NTT NGN での IPv6」(p.437) — 日本固有の話が全部ここにある",
        fsize=13,
    )

    # ================= 第0部: 全体像(まず地図を配る) =================
    section_slide(prs, "第0部", "全体像 — まず地図を配る",
                  "細かい話に入る前に、これから何を何の順で見るのかを一枚で掴む")

    bullet_slide(
        prs,
        "この資料の歩き方(全体 → 詳細 の順に降ります)",
        [
            (0, "第0部 全体像 — いま何の話をしているのか", "b"),
            (1, "IPv6 の話は「アドレスそのもの」「配り方」「IPv4 の運び方」の 3 層に分かれる", ""),
            (0, "第1部 IPv6 そのもの(抽象)— どこの国でも同じ話", "b"),
            (1, "アドレスの書き方 / 種類 / 決まり方。ここは教科書どおり", ""),
            (0, "第2部 日本の回線(具体)— ここからは日本だけの話", "b"),
            (1, "NTT の網 + VNE + ISP。v6プラス / transix などが「どれ」なのかを表で押さえる", ""),
            (0, "第3部 実習(第1回〜第3回)— 手を動かして原理を見る", "b"),
            (1, "MAP-E のアドレス計算、ポートの割り当て、MTU の引き算。ラボで壊して直す", ""),
            (0, "第4部 検証結果(実測)— 本当にそうなるのか", "b"),
            (1, "ラボと実機で測った値。「たぶん」と「測った」を必ず分けて書いてあります", ""),
            (0, "迷ったら第0部に戻ること。地図の上のどこにいるか分かれば必ず戻れます", "w"),
        ],
    )

    bullet_slide(
        prs,
        "3 行でいうと",
        [
            (0, "① PPPoE が遅い", "b"),
            (1, "夜になると遅い。原因は網終端装置(NTE)の輻輳で、回線速度の問題ではない", ""),
            (0, "② だから IPoE に変える", "b"),
            (1, "IPoE = 網終端装置を通らず、網から直接 IPv6 が降ってくる方式", ""),
            (0, "③ でも社内は IPv4 のままなので、IPv4 を IPv6 の上に載せて運ぶ", "b"),
            (1, "その運び方が MAP-E と DS-Lite。この資料の 8 割はここの話です", ""),
            (0, "", ""),
            (0, "ここだけ持ち帰れば十分です。残りは全部この 3 行の詳細でしかありません", "w"),
        ],
    )

    bullet_slide(
        prs,
        "なぜ IPv6 に変えるのか — 建前と本音を分ける",
        [
            (0, "教科書に書いてある理由(建前)", "b"),
            (1, "IPv4 アドレスが枯渇した。だから IPv6 へ移行する", ""),
            (1, "正しいが、これは「きっかけ」であって、お客様が今ハンコを押す理由ではない", ""),
            (0, "現場でお金が動く理由(本音)", "b"),
            (1, "PPPoE が遅いから。「夜だけ遅い」「Web 会議が切れる」で相談が来る", ""),
            (1, "IPv6 を使いたいのではなく、速度が欲しくて IPv6 の道を通ることになる", ""),
            (0, "この区別が効く場面", "b"),
            (1, "「社内も IPv6 化しますか?」→ 普通はしません。LAN は IPv4 のまま", ""),
            (1, "アドレス再設計・FW 見直し・アプリ検証が発生する。要件でない限りやらない", ""),
            (1, "IPv6 は「外を速く通るための道」。社内を作り替える話ではない", ""),
            (0, "※ ただし「将来も IPv6 化しない」とは言わないこと。今回やらないだけです", "w"),
        ],
        footer="教科書: 第18章「IPv4 アドレスの枯渇」/ 付録 A「NTT NGN での IPv6」(p.437)",
    )

    diagram_slide(
        prs,
        "IPv6 の話は 3 層に分かれる(混ぜると分からなくなる)",
        [
            ("① アドレスそのもの", "書き方・種類\n128ビット / fe80:: / /64"),
            ("② アドレスの配り方", "RA(SLAAC)/ DHCPv6-PD\n誰が何を配るか"),
            ("③ IPv4 の運び方", "MAP-E / DS-Lite\nIPv4 を IPv6 に載せる"),
        ],
        note=[
            "トラブルの切り分けも、この 3 層のどこで壊れているかを決めるところから始まります。",
            "「IPv6 が降りてこない」= ② の話。「IPv6 は降りたが IPv4 が通らない」= ③ の話。",
            "※ ①② は世界共通の技術。③ だけが日本固有の事情が濃い(共有 IPv4 とポート制限)。",
        ],
    )

    # ================= 第1部: IPv6 そのもの =================
    section_slide(prs, "第1部", "IPv6 そのもの — どこの国でも同じ話",
                  "アドレスの書き方・種類・決まり方。暗記より「見たことがある」を作る")

    table_slide(
        prs,
        "IPv4 と IPv6 の違い(現場で効くところだけ)",
        ["", "IPv4", "IPv6", "現場でどう効くか"],
        [
            ["長さ", "32 ビット", "128 ビット", "枯渇しない。ただし覚えられない → コピペ前提で作業する"],
            ["表記", "192.168.1.1", "2001:db8::1", ":: は 0 の連続の省略。1 か所しか使えない"],
            ["1 サブネット", "任意(/24 等)", "/64 が原則", "/64 より小さく切ると SLAAC が動かない"],
            ["1 IF のアドレス数", "普通 1 個", "複数が当たり前", "リンクローカル + グローバルが同時に付く"],
            ["ブロードキャスト", "あり", "無い", "代わりにマルチキャスト。ARP も無い(NDP になる)"],
            ["アドレス解決", "ARP", "NDP(ICMPv6)", "ICMPv6 を止めると通信が死ぬ。FW 設計の最頻出事故"],
            ["自動設定", "DHCP", "RA / DHCPv6", "ルータが自分から「私がゲートウェイだ」と言いに来る"],
            ["フラグメント", "途中のルータも分割", "送信元だけ", "経路 MTU の問題が表に出やすい(後述の 1460)"],
        ],
        [1.7, 2.1, 2.3, 5.7],
        footer="教科書: 2.10「IPv6 と IPv4 の違い」(p.51)",
        fsize=11,
    )

    bullet_slide(
        prs,
        "アドレスの読み方と、なぜ /64 なのか",
        [
            (0, "書き方の 3 つのルール", "b"),
            (1, "各ブロック先頭の 0 は省ける: 2001:0db8:0000:0001 → 2001:db8:0:1", ""),
            (1, "0 が続くところは :: で 1 回だけ省ける: 2001:db8:0:0:0:0:0:1 → 2001:db8::1", ""),
            (1, "URL に書くときは [] で囲む(忘れると動きません)", ""),
            (1, "http://[2001:db8::80]/", "c"),
            (0, "/64 が原則である理由", "b"),
            (1, "後ろの 64 ビットをインタフェース ID に使う、と決まっているため", ""),
            (1, "MAC から作る(EUI-64)か、ランダムに作る。つまり下位は端末が勝手に決める", ""),
            (1, "→ /64 より小さく切ると SLAAC が動かない。IPv4 の感覚をここで捨てる", "w"),
            (0, "実務での見え方(ここが後で効きます)", "b"),
            (1, "ひかり電話なし = /64 が 1 本だけ来る(RA 方式)", ""),
            (1, "ひかり電話あり = /56 が委譲される = /64 が 256 本使える(PD 方式)", ""),
            (1, "この違いが効くのは「LAN をいくつのセグメントに切れるか」と「HGW とどちらで終端するか」", ""),
            (1, "※ MAP-E は /64 (RA 方式) でも組めます。ラボでも RA 方式で 198.51.100.20 / PSID 3 を導出しました", "w"),
        ],
        footer="教科書: 第3章「IPv6 アドレス」/ 第7章「SLAAC」",
    )

    table_slide(
        prs,
        "アドレスの種類 — 1 つの NIC に何本も付く",
        ["種類", "見た目", "届く範囲", "誰が決める", "現場でどこで見るか"],
        [
            ["リンクローカル(LLA)", "fe80::…", "同じリンク内だけ", "自動(必ず付く)", "NDP と RA の送信元。ルーティングの隣接もこれ"],
            ["グローバル(GUA)", "2000::/3", "インターネット全体", "RA か DHCPv6", "実際の通信に使うアドレス"],
            ["ユニークローカル(ULA)", "fd00::/8", "組織内だけ", "自分で決める", "規格上は fc00::/7。OpenWrt が既定で作る (事故の元)"],
            ["マルチキャスト", "ff00::/8", "グループ宛", "規格で決まっている", "ff02::1 全ノード / ff02::2 全ルータ"],
            ["エニーキャスト", "GUA と同じ", "一番近い 1 台", "設計で決める", "DNS など。見た目では区別できない"],
            ["未指定", "::", "—", "—", "DAD の送信元。「まだ決まっていない」の意味"],
            ["ループバック", "::1", "自分だけ", "—", "IPv4 の 127.0.0.1 に相当"],
        ],
        [2.6, 1.6, 2.0, 2.1, 3.9],
        footer="資料や検証に使う「ドキュメント用」は 2001:db8::/32。実在アドレスを例に使わないこと",
        fsize=11,
    )

    bullet_slide(
        prs,
        "fe80:: が必ず付いている理由 — IPv4 と一番違うところ",
        [
            (0, "IPv6 は「アドレスをもらう前に通信する」必要がある", "b"),
            (1, "アドレスをくれるのはルータ。でもルータと話すにはアドレスが要る → 鶏と卵", ""),
            (1, "だから NIC が上がった瞬間に自分だけで作れるアドレスが要る。それが fe80::", ""),
            (0, "だから、こうなっている", "b"),
            (1, "RA も DHCPv6 も送信元は fe80::。グローバルは後から付く", ""),
            (1, "既定ゲートウェイも fe80:: で指す。IPv4 のように 192.168.1.1 とは書かない", ""),
            (1, "OSPFv3 / BGP の隣接も fe80::。「グローバルが無い = 壊れている」ではない", ""),
            (0, "現場で効く場面", "b"),
            (1, "「グローバル IPv6 が無いのに ping6 が通る」→ 正常。LLA で通っています", ""),
            (1, "このラボでも、VM に管理 IP を振らず fe80:: だけで SSH して構築しました", ""),
            (1, "※ fe80:: は どのインタフェース経由か を書かないと使えない(%eth0 / %vmbr0)", "w"),
        ],
        footer="教科書: 第5章「ICMPv6」/ 第6章「近隣探索プロトコル(NDP)」",
    )

    branch_slide(
        prs,
        "アドレスはどうやって決まるのか",
        ("ルータから RA が飛んでくる", "ICMPv6。定期送信 + 要求への応答"),
        ("SLAAC(自動設定)", "RA にプレフィックスが入っている\n端末が下位 64 ビットを自分で作る"),
        ("DHCPv6", "RA は「DHCPv6 を使え」と言うだけ\nアドレスは DHCPv6 サーバが配る"),
        note=[
            "RA の中の 2 つのフラグ(M / O)で、どちらを使うかが決まります。",
            "M=0 → DHCPv6 でアドレスはもらわない。M=1 → DHCPv6 でも アドレスをもらう。O=1 → DNS 等だけ DHCPv6 から。",
            "※ SLAAC するかどうかを決めるのは M ではなく、RA の中のプレフィックス情報にある A フラグです。",
            "  M=1 でも A=1 なら SLAAC も併走し、端末にアドレスが 2 本以上付きます。",
            "※ M/O はルータからの推奨でしかなく、従うかは OS 次第 (Android は DHCPv6 を使いません)。",
            "※ ルータがもらう /56 は DHCPv6-PD という別枠。アドレスではなくプレフィックスの委譲です。",
            "※ 日本の小規模案件では「ひかり電話あり = PD、なし = RA だけ」で実務はほぼ回ります。",
        ],
    )

    # ================= 第2部: 日本の回線 =================
    section_slide(prs, "第2部", "日本の回線 — ここからは日本だけの話",
                  "登場人物と、v6プラス / transix などが結局「どれ」なのかを表で押さえる")

    diagram_slide(
        prs,
        "登場人物 — 「ISP」と一言で言っても 3 者いる",
        [
            ("お客様の CPE", "ルータ\n892FJ / RTX / OpenWrt"),
            ("NTT の網(NGN)", "フレッツ光\nここまでは NTT"),
            ("VNE", "IPv6 を実際に流す事業者\nJPIX / MF / ドコモビジネス 等"),
            ("ISP", "契約先・課金\nVNE とは限らない"),
        ],
        note=[
            "「ISP を変えたのに速度が変わらない」の正体はここ。速度を決めるのは VNE 側です。",
            "お客様は ISP としか契約していないつもりでも、裏で VNE が決まっており、",
            "その VNE が MAP-E か DS-Lite かを決めています。こちらでは選べません。",
            "※ 提案・見積の前に「ISP はどこか」ではなく「VNE はどこか」を聞くこと。",
        ],
    )

    branch_slide(
        prs,
        "分かれ道はまず 1 段目 — PPPoE か IPoE か",
        ("フレッツ光の契約", "どちらで IPv4/IPv6 を通すか"),
        ("PPPoE(従来)", "ID/パスワードで認証\n網終端装置を通る = 夜が遅い\nIPv4 がそのまま流れる"),
        ("IPoE", "認証なし。回線に紐づく\n網終端装置を通らない = 速い\nIPv4 は別途トンネルが要る"),
        note=[
            "1 段目はここ。速度の話はこの分岐で決まります。",
            "2 段目は「IPoE を選んだあと、IPv4 をどう運ぶか」= MAP-E か DS-Lite か(次ページ)。",
            "※ 両方同時に使えます。IPv6 は IPoE、IPv4 は PPPoE、という併用は実際によくある構成です。",
            "※ 切替当日は PPPoE を落としてから IPoE。セッションが残ると原因不明の不通になります。",
        ],
    )
    branch_slide(
        prs,
        "2 段目 — IPv4 をどう運ぶか(ここは選べません)",
        ("IPoE にした。ではIPv4 は?", "IPv6 の上に載せて運ぶしかない"),
        ("MAP-E", "CPE が IPv4 を IPv6 で包む\n共有 IPv4 + 使えるポートが決まる\nCPE 側で計算が要る"),
        ("DS-Lite", "CPE はただ包むだけ\nNAT は網側(AFTR)がやる\nCPE の負荷が軽い"),
        note=[
            "どちらになるかは VNE が決めます。お客様も我々も選べません。",
            "違いが一番出るのは「ポート開放できますか?」への答え(第3部で詳しく)。",
            "※ MAP-E: 使えるポートが数百個に限られる。指定のポートは開けられないことが多い。",
            "※ DS-Lite: そもそも利用者側に開ける手段がない。固定 IP オプションを買う話になる。",
        ],
    )

    table_slide(
        prs,
        "実際のサービスはどれ? — 名前から方式を引く表",
        ["サービス名", "VNE / 提供", "方式", "ポート開放"],
        [
            ["v6プラス", "JPIX(旧 JPNE)", "MAP-E", "限定的(割当ポートのみ)"],
            ["OCN バーチャルコネクト", "NTT ドコモビジネス(旧 NTT Com)", "MAP-E", "限定的(割当ポートのみ)"],
            ["BIGLOBE IPv6 オプション", "BIGLOBE 系", "MAP-E", "限定的(割当ポートのみ)"],
            ["transix", "インターネットマルチフィード", "DS-Lite", "不可(固定 IP 契約が別途必要)"],
            ["クロスパス(Xpass)", "アルテリア・ネットワークス", "DS-Lite", "不可(固定 IP 契約が別途必要)"],
            ["v6 コネクト", "朝日ネット", "DS-Lite", "不可(固定 IP 契約が別途必要)"],
            ["ソフトバンク光 IPv6 高速ハイブリッド", "BBIX", "独自方式", "光BBユニット前提。本資料の対象外"],
        ],
        [3.6, 3.6, 1.6, 3.2],
        footer="出典: ヤマハ RTpro「動作確認済み IPv6 IPoE + IPv4 over IPv6 接続サービス」ほか各社資料。案件では必ず最新版を確認すること",
        fsize=12,
    )

    bullet_slide(
        prs,
        "この表の使い方と、使うときの注意",
        [
            (0, "ヒアリングでこの順に聞く", "b"),
            (1, "① ひかり電話はありますか → /64 か /56 かが決まる(= PD 方式かどうか)", ""),
            (1, "② 契約は ISP どこですか → その裏の VNE を調べる", ""),
            (1, "③ VNE が決まれば方式が決まる → ポート開放の可否が言える", ""),
            (0, "注意 1: 同じ ISP でもプランで VNE が違うことがある", "w"),
            (1, "「◯◯光」は再販が多く、同じ名前でも中身が別 VNE ということが起きる", ""),
            (1, "机上で決めつけず、開通後に実測で確かめること(第4部の run-checks)", ""),
            (0, "注意 2: 事業者名が変わっている", "w"),
            (1, "JPNE は JPIX に、NTT コミュニケーションズは NTT ドコモビジネスに", ""),
            (1, "古い資料と名前が食い違うので、社内 wiki の記述は更新すること", ""),
            (0, "注意 3: 方式が分かっても「動く」とは限らない", "w"),
            (1, "CPE 側が対応している必要があります。次のページが実測結果です", ""),
        ],
    )

    # ---------------- 第1回 ----------------
    section_slide(prs, "第1回", "なんで IPv6 に変えなきゃいけないの?",
                  "登場人物を絵で描ける。fe80:: とグローバルの違いが言える")

    bullet_slide(
        prs,
        "1-0. 登場人物と略語(まずこれを覚える)",
        [
            (0, "お客様PC ─ CPE ─ [ NGN ] ─ VNE ─ ISP ─ インターネット", "b"),
            (1, "ひかり電話ありなら CPE の位置に HGW が入る", ""),
            (0, "CPE … Customer Premises Equipment = お客様側に置くルータ", ""),
            (0, "HGW … ひかり電話ルータ。機種により IPoE を自分で終端してしまう", "w"),
            (0, "NGN … NTT東西の IP 網。ISP とは別の会社の設備", ""),
            (0, "VNE … NGN 上の IPv6 (IPoE) 接続を ISP に卸す事業者。多くは IPv4 over IPv6 もセットで提供", ""),
            (1, "JPIX (旧 JPNE) / インターネットマルチフィード / NTT ドコモビジネス (旧 NTT Com) 等", ""),
            (0, "BRAS … PPPoE を終端する ISP 側の装置", ""),
            (0, "NTE … 網終端装置。PPPoE が必ず通る箱で、ここが夕方混む", ""),
            (0, "ラボの VM との対応", "b"),
            (1, "NGN-SIM = NGN / BRAS = BRAS / VNE = VNE (BR と AFTR 兼用) / INET-SIM = インターネット", ""),
            (1, "OpenWrt-CE または実機 = CPE。ラボに NTE は無く、BRAS が代役", ""),
        ],
        footer="略語の全量は study-guide.md §1.4。ここが分からないと以降ずっと分からない",
    )

    bullet_slide(
        prs,
        "1-1. まず事実: IPv4 アドレスは在庫が尽きた",
        [
            (0, "IPv4 は約 43 億個。世界人口より少ない。スマホ 1 人 1 台で足りない", ""),
            (0, "だから 3 つの対策が同時に走っている", "b"),
            (1, "① NAT で共有する → 家庭内では昔からやっている", ""),
            (1, "② アドレスを売買・移転する → 実際に市場がある", ""),
            (1, "③ IPv6 に移る → 本命", ""),
            (0, "IPoE への切替は ③ の話。「速いから」は結果であって理由ではない", "w"),
            (0, "教科書: 18.1 在庫枯渇はどのような問題なのか (p.335) / 18.2 影響 (p.336)", ""),
        ],
    )

    bullet_slide(
        prs,
        "1-2. IPv6 アドレスの書き方 (ここで演習)",
        [
            (0, "128 ビット。16 進 4 桁を 8 ブロック、コロン区切り", ""),
            (1, "2001:0db8:1014:0300:0000:0000:0000:0001", "c"),
            (0, "省略ルールは 2 つだけ", "b"),
            (1, "① 各ブロック先頭の 0 は省略できる → 0db8 は db8、0300 は 300", ""),
            (1, "② 連続する 0 のブロックは :: に潰せる。ただし 1 箇所だけ", ""),
            (1, "2001:db8:1014:300::1", "c"),
            (0, "演習: 次を省略形にせよ (紙に書く。答え合わせまで含めて 20 分)", "b"),
            (1, "fe80:0000:0000:0000:0200:5eff:fe00:5301", "c"),
            (1, "2001:0db8:0000:0000:0000:0000:8888:0001", "c"),
            (1, "2001:0db8:100a:0500:0000:0000:0000:0000", "c"),
            (0, "教科書: 2.2 IPv6 アドレスのテキスト表記 (p.38)", ""),
        ],
        footer="※ ここを飛ばすと第2回以降アドレスが読めないまま進む。時間を削らないこと",
    )

    diagram_slide(
        prs,
        "1-3. IPv6 と IPv4 は「別のインターネット」",
        [("IPv4 の世界", "203.0.113.x"), ("互換性なし", "直接は話せない"), ("IPv6 の世界", "2001:db8::x")],
        [
            "IPv4 と IPv6 の間に自動変換はない。だから両方を同時に使う (デュアルスタック) か、",
            "どちらかを他方の上で運ぶ (共存技術) しかない。IPoE がやっているのは後者。",
            "※ 「IPv6 にしたら IPv4 が使えなくなる」は誤り。IPv6 のトンネルの中を通って使える。",
            "教科書: 2.10 IPv6 と IPv4 の違い (p.51) / 2.12 同時に使う (p.58) / 2.13 共存技術 (p.61)",
        ],
    )

    bullet_slide(
        prs,
        "1-4. 演習 1-A: CPE の WAN に何が付いているか  (手順は実測済み)",
        [
            (0, "手順", "b"),
            (1, "1. CPE (OpenWrt-CE または実機。クライアント VM ではない) で PPPoE 接続", ""),
            (1, "   認証: user1@isp-a.example / pass1 / サービス名 lab-isp", "c"),
            (1, "2. CPE の WAN 側で  ip -6 addr show  (実機なら show ipv6 interface brief)", "c"),
            (1, "3. 出てきたアドレスを 1 個ずつ「これは何か」に分類する", ""),
            (0, "ラボの仕掛け — ここが今日の山場", "b"),
            (1, "ラボの BRAS は ipv6=deny。PPPoE セッション自体では IPv6 は降りない", ""),
            (1, "(= IPv6 非対応の PPPoE 契約の再現)", ""),
            (1, "しかし CPE の WAN は NGN アクセス網の L2 に直結していて、NGN-SIM の radvd は", ""),
            (1, "ra/pd どちらのモードでも常時 RA を流している", ""),
            (0, "つまり PPPoE とは無関係に、NGN 由来のグローバル IPv6 が来ている可能性がある", "w"),
            (1, "これが 付録 A.1 フォールバック問題 / A.2 マルチプレフィックス問題 の正体", ""),
            (1, "「ISP の契約とは別に NGN の IPv6 が勝手に来ている」を最初に体感させる", ""),
        ],
        footer="手順はラボで実走済み。当日の受講者環境で GUA がどのプレフィックスで付くかは、その場で study-guide に記録すること",
    )

    bullet_slide(
        prs,
        "1-5. リンクローカルとグローバル / 宿題",
        [
            (0, "fe80:: で始まる = リンクローカル。同じリンク内専用。消すと全部死ぬ", "b"),
            (1, "IPv6 のルーティングも近隣探索もこれを使って動く。RA の送信元も fe80::", ""),
            (1, "IPv4 の 169.254.x.x とは役割が違う (あっちは DHCP 失敗時の敗戦処理)", ""),
            (0, "2001: で始まる = グローバル。外と話せる", "b"),
            (0, "Temporary (プライバシー拡張) について", "b"),
            (1, "lab-client (Ubuntu Server) では既定で生成されない", ""),
            (1, "見せたい場合: sudo sysctl -w net.ipv6.conf.<if>.use_tempaddr=2", "c"),
            (0, "宿題: 自宅の PC で IPv6 アドレスを見て、種類を数えてくる", "b"),
            (1, "Windows: ipconfig /all   /   Mac・Linux: ip -6 addr", "c"),
            (1, "自宅が IPv6 非対応・管理画面に入れない場合はスマホのテザリングで代替可", ""),
            (0, "教科書: 3.5 リンクローカル (p.74) / 2.9 複数のアドレスが付く (p.48) / 7.5 Temporary (p.166)", ""),
        ],
    )

    # ---------------- 第2回 ----------------
    section_slide(prs, "第2回", "アドレスはどうやって降ってくるのか",
                  "ひかり電話の有無と「どこで終端するか」の 2 段で方式を決められる")

    bullet_slide(
        prs,
        "2-1. IPoE とは何か (用語の整理から)",
        [
            (0, "IPoE = IP over Ethernet。PPP を挟まずに Ethernet に直接 IP を載せる方式", "b"),
            (1, "NTT の商品名ではない。接続方式の名前", "w"),
            (0, "PPPoE との違い", "b"),
            (1, "PPPoE: ID/パスワードでセッションを張る → 網終端装置 (NTE) を必ず通る", ""),
            (1, "IPoE: つないだらアドレスが降ってくる → NTE を通らない", ""),
            (0, "「IPoE は速い」の正体", "b"),
            (1, "プロトコルが速いのではなく、混雑する箱 (NTE) を通らないから", ""),
            (1, "※ VNE 側が混んでいれば普通に遅い。「速くなります」と客に約束しない", "w"),
            (0, "教科書: 付録 A.3 IPv6 PPPoE と IPv6 IPoE (p.443) ← 付箋を貼る", ""),
        ],
    )

    bullet_slide(
        prs,
        "2-2. 近隣探索プロトコル — IPv6 の心臓部",
        [
            (0, "IPv6 に ARP はない。代わりに ICMPv6 の 4 兄弟が働く (ほかに Redirect もある)", "b"),
            (1, "RS (Router Solicitation): 「ルータいますか?」", ""),
            (1, "RA (Router Advertisement): 「ここにいる。プレフィックスはこれ」", ""),
            (1, "NS (Neighbor Solicitation): 「このアドレスの MAC 教えて」= ARP 相当", ""),
            (1, "NA (Neighbor Advertisement): 「私です」", ""),
            (0, "全部 ICMPv6。だから ICMPv6 を全 deny すると IPv6 は動かない", "w"),
            (1, "教科書 14.9「ICMPv6 を無条件にすべてフィルタリングすべきではない」(p.279)", ""),
            (1, "対立意見: では全許可か? それも違う。RFC 4890 相当の選別が必要", ""),
            (0, "教科書: 6.1 機能と利用するメッセージ (p.128) / 6.2 ルータとプレフィックス情報の発見 (p.129)", ""),
        ],
    )

    bullet_slide(
        prs,
        "2-3. 2 つの方式 — RA 方式と PD 方式",
        [
            (0, "RA 方式 (SLAAC)", "b"),
            (1, "ルータ広告のプレフィックス情報から端末が自分でアドレスを作る", ""),
            (1, "降りてくるのは /64 が 1 個。LAN を複数セグメントには分けられない", ""),
            (0, "PD 方式 (DHCPv6-PD)", "b"),
            (1, "DHCPv6 で「プレフィックスそのもの」が委任される", ""),
            (1, "/56 なら /64 が 256 個作れる。これが PD 方式の価値", ""),
            (0, "IPv4 との感覚の違い", "b"),
            (1, "IPv4 は「アドレスを 1 個もらう」、IPv6 は「プレフィックスをもらって中を自分で切る」", ""),
            (1, "SLAAC を使う限りホスト部は 64 ビット。だから /64 より細かくは切らない", ""),
            (0, "教科書: 7.1 SLAAC の流れ (p.160) / 8.1 DHCP と DHCPv6 の違い (p.180) / 8.7 PD (p.201)", ""),
        ],
    )

    branch_slide(
        prs,
        "2-4. 方式判定は 2 段。1 段で決めると事故る",
        ("① ひかり電話は?", "ヒアリングシート #1"),
        ("なし → RA 方式 (/64)", "配下ルータで終端"),
        ("あり → HGW が来る → ② へ", "HGW 型番の確認が必須 (#2)"),
        [
            "② 終端するのはどっち?   HGW で終端 → 配下ルータは PD を要求しない (ブリッジ/ルータモード)",
            "                        配下ルータで終端 → HGW 側の IPoE を止める",
            "※ 「ひかり電話あり = 配下ルータに PD を設定」は誤り。HGW が PD を取るので二重終端 (R7) になる。",
            "教科書: 付録 A.3 (p.443) / 付録 B ND Proxy (p.455) / 13.2 マルチプレフィックスの問題 (p.251)",
        ],
    )

    table_slide(
        prs,
        "2-5. NGN と VNE はセットで切り替える  (実測済み)",
        ["NGN-SIM モード", "ユーザプレフィックス", "共有 IPv4 / PSID", "VNE 側 (BR) の実行方法"],
        [
            ["pd (ひかり電話あり)", "2001:db8:100a:500::/56", "198.51.100.10 / 5", "sudo ./setup-map-br.sh"],
            ["ra (ひかり電話なし)", "2001:db8:1014:300::/64", "198.51.100.20 / 3",
             "!環境変数を付けて再実行 (下記)"],
        ],
        [2.6, 3.5, 2.5, 3.6],
        footer="ra に切り替えたら → sudo CE_MAP_ADDR=2001:db8:1014:300:0:c633:6414:3 CE_SHARED_V4=198.51.100.20 ./setup-map-br.sh\n※ NGN 側だけ切り替えて BR を再実行しないと IPv4 が全断する (build.md §3)",
        fsize=12,
    )

    bullet_slide(
        prs,
        "2-6. 演習 2-A: 4 兄弟をキャプチャする  (実測済み)",
        [
            (0, "手順 (root 必須。sudo を付ける)", "b"),
            (1, "sudo ./lab/ipoe/ngn/setup-ngn.sh ra      # ひかり電話なし相当", "c"),
            (1, "↑ 切り替えたら VNE でも必ず setup-map-br.sh を再実行 (前ページの表)", "w"),
            (1, "sudo tcpdump -nn -i <ACCESS_IF> 'icmp6 or udp port 546 or udp port 547'", "c"),
            (0, "見てほしいこと", "b"),
            (1, "RS → RA が飛ぶか。PD なら Solicit → Advertise → Request → Reply の 4 往復", ""),
            (1, "ip -6 route と CPE の LAN 側アドレスを比べる", "c"),
            (1, "RA 方式だと LAN に配れる /64 が無い ← だから PD が必要 (ラボは odhcpd relay で回避)", ""),
            (0, "先に済ませておくこと(当日ハマる)", "w"),
            (1, "インターフェース名を確定しておく (detect-ifs.sh が MAC から解決する)", ""),
            (1, "OpenWrt の wan は論理名。tcpdump には実体名 (eth1 等) が必要", ""),
            (1, "OpenWrt に tcpdump は既定で入っていない。ラボ内にインターネットは無い", ""),
            (1, "setup-ngn.sh は毎回 apt-get する。管理経路が出られないと演習中に落ちる", ""),
        ],
    )

    bullet_slide(
        prs,
        "2-7. 演習 2-B: DUID の罠",
        [
            (0, "OpenWrt は DUID-LL を送るので DROP を素通りする (サイクル 5 実測)", "w"),
            (1, "README.md / test-matrix.md §3 Phase 1 手順 5", ""),
            (0, "見せる手順", "b"),
            (1, "1. kea-dhcp6.conf の client-classes (DROP クラス) のコメントを外す", ""),
            (1, "2. Kea を再起動 → CPE から Solicit", ""),
            (1, "3. エラーではなく完全な無応答になることを見せる", ""),
            (1, "4. 終わったら戻す", ""),
            (0, "用語", "b"),
            (1, "DUID-LL  = MAC アドレスベース (type 0003)", ""),
            (1, "DUID-LLT = MAC + 生成時刻", ""),
            (0, "NGN は LL しか受けないという報告が多い。ただし公開仕様では確認できていない", "w"),
            (1, "断定せず「そういう報告が多いので DUID-LL で試す」と伝えるのが安全", ""),
            (0, "教科書: 8.4 DUID (p.191)。知らないと「なぜか PD だけ来ない」で数時間溶ける", ""),
        ],
    )

    # ---------------- 第3回 ----------------
    section_slide(prs, "第3回", "IPv4 を IPv6 の上で運ぶ",
                  "「ポート開放したい」に、方式を聞いて可否を即答できる")

    bullet_slide(
        prs,
        "3-0. この回は分量オーバー。絞る",
        [
            (0, "読み上げない (資料配布のみ)", "w"),
            (1, "共存技術の分類 (トンネル / 変換 / プロキシ) … 21.1 (p.381)", ""),
            (1, "6to4 / Teredo / ISATAP / 6rd / 4rd の系譜 … 第22章 (p.385)", ""),
            (1, "SIIT / NAT64+DNS64 / 464XLAT … 第23章 (p.407)。モバイル系で固定回線では少ない", ""),
            (0, "この 90 分でやること", "b"),
            (1, "なぜ IPv6 プレフィックスの中に IPv4 が埋まっているのか (第2回からの橋)", ""),
            (1, "MAP-E: EA bits と PSID を手で計算する", ""),
            (1, "DS-Lite: NAT を網側に置く。だからポート開放できない", ""),
            (1, "MTU の引き算", ""),
            (1, "演習 3-A: ポートを数える / R3 / R4 / R5", ""),
            (0, "教科書で読むのは 24 章だけでよい (p.421-)", "b"),
        ],
    )

    bullet_slide(
        prs,
        "3-1. サービス名と方式名を混同しない (最重要の実務知識)",
        [
            (0, "これは全部サービス名 (商品名)。教科書には出てこない", "w"),
            (1, "v6 プラス / transix / クロスパス / IPv6 オプション / v6 コネクト …", ""),
            (0, "方式名はこの 2 つ (固定回線の日本の IPoE では実質これだけ)", "b"),
            (1, "MAP-E … CE (お客様のルータ) 側で NAT する。ポートが制限される", ""),
            (1, "DS-Lite … 網側 (AFTR) で NAT する。利用者にポート開放の手段がない", ""),
            (0, "案件で最初に確認すること", "b"),
            (1, "「その商品の中身は MAP-E ですか DS-Lite ですか」", ""),
            (1, "サービス名だけで設定を始めると、CPE のテンプレを間違える", "w"),
            (0, "実サービスとの対応表: 第2部の早見表 / study-guide.md §1.35", ""),
            (0, "※ ソフトバンク光は例外です", "w"),
            (1, "BBIX の独自方式 (IPv6 高速ハイブリッド。4rd/SAM 系) で、光BBユニット前提", ""),
            (1, "MAP-E でも DS-Lite でもないので、本資料の設計対象外として扱ってください", ""),
        ],
    )

    bullet_slide(
        prs,
        "3-2. MAP-E のポート数はこう決まる(計算する)",
        [
            (0, "ラボのルール (build.md §3)", "b"),
            (1, "rule-ipv6-prefix 2001:db8:1000::/40 / rule-ipv4-prefix 198.51.100.0/24", "c"),
            (1, "ea-len 16 / psid-offset 4", "c"),
            (0, "式", "b"),
            (1, "psid-len = ea-len − (32 − v4プレフィックス長) = 16 − (32 − 24) = 8", "c"),
            (1, "1ブロックのポート数 = 2^(16 − offset − psid-len) = 2^4 = 16", "c"),
            (1, "ブロック数 = 2^offset − 1 = 15   ← ブロック 0 (ポート 0〜4095) は使わない", "c"),
            (1, "使えるポート数 = 16 × 15 = 240", "c"),
            (0, "実サービスの実績値 (research-notes.md §2)", "b"),
            (1, "v6プラス (JPIX) = 240 ポート / OCN バーチャルコネクト (NTT ドコモビジネス) = 1008 ポート", ""),
            (0, "240 は少ない。ブラウザは 1 タブで数十本張る", "w"),
            (1, "台数の多い拠点や監視系があると足りない。同時セッション数を実測してから方式を決める", ""),
        ],
    )

    bullet_slide(
        prs,
        "3-3. PD 方式 (/56) から実際に導出してみる",
        [
            (0, "与えられたプレフィックス: 2001:db8:100a:0500::/56", "c"),
            (0, "① EA bits = プレフィックスの bit40〜55 を取り出す → 0x0a, 0x05", "b"),
            (0, "② 共有 IPv4 = 198.51.100.0/24 の下位 8bit に 0x0a を入れる → 198.51.100.10", "b"),
            (0, "③ PSID = 残り 8bit → 5", "b"),
            (0, "④ ポート集合", "b"),
            (1, "0x1050-0x105F (4176-4191)", "c"),
            (1, "0x2050-0x205F, 0x3050-0x305F, … 0xF050-0xF05F", "c"),
            (1, "16 個 × 15 ブロック = 240。連続していない", "w"),
            (0, "RA 方式 (2001:db8:1014:0300::/64) なら → 共有 IPv4 198.51.100.20 / PSID 3", ""),
            (0, "「4176〜8000 の範囲」のように連続だと思って NAT を書くと動かない → R12", "w"),
            (0, "教科書: 24.3 A+P (p.423) / 24.4 MAP-E、MAP-T、4rd (p.425)", ""),
        ],
    )

    diagram_slide(
        prs,
        "3-4. DS-Lite — NAT が網側にある。だから触れない",
        [("お客様 PC", "192.168.x.x"), ("CPE (B4)", "NAT しない"), ("AFTR", "ここで NAT"), ("インターネット", "共有グローバル")],
        [
            "MAP-E は CPE (CE) が NAT する → 自分のポート集合の中なら開放できる。",
            "DS-Lite は AFTR (事業者設備) が NAT する → 利用者からそれを操作する手段が提供されていない。",
            "※ 「ポート開放できますか?」→ DS-Lite なら「できません」。代替は固定 IP 系サービスか PPPoE 併用。",
            "教科書: 24.1 DS-Lite (p.421) / 19.3 NAT 機器に要求される挙動 (p.356) / 19.4 課題 (p.362)",
        ],
    )

    bullet_slide(
        prs,
        "3-5. MTU の引き算と、直し方",
        [
            (0, "連続した引き算ではありません。PPPoE 系と IPoE トンネル系の 2 系統です", "b"),
            (0, "【PPPoE の場合】 1500 − PPPoE ヘッダ 8 = 1492 (MRU 上限で実質 1454 が多い)", "c"),
            (0, "【MAP-E / DS-Lite の場合】 1500 − IPv6 ヘッダ 40 = 1460", "c"),
            (0, "     さらに encaplimit (Dst Options) が付くと −8 → 1452", "c"),
            (0, "     ip6tnl の既定でこれが付く。片方向だけ切れる原因になる (R11)", "c"),
            (0, "症状の出方が独特なので覚える", "w"),
            (1, "小さいパケット (HTTP GET, ping, DNS) は通る → 「つながってる」と見える", ""),
            (1, "大きいパケット (ファイル送信、添付) だけ落ちる", ""),
            (0, "現場の一手 — 症状を知っているだけでは直せない", "b"),
            (1, "TCP MSS clamp: OpenWrt は firewall の mtu_fix、Cisco は ip tcp adjust-mss 1420", "c"),
            (1, "CPE の WAN MTU を手動設定: OpenWrt なら option mtu '1460'", "c"),
            (1, "トンネルの encaplimit none (ラボの VNE スクリプトは設定済み)", ""),
            (1, "経路で ICMP (v6 Packet Too Big / v4 Fragmentation-Needed) が落ちていないか確認", ""),
            (0, "教科書: 第10章 (p.215) / 9.2 IPv6 フラグメントヘッダ (p.208) / 14.9 (p.279)", ""),
        ],
    )

    bullet_slide(
        prs,
        "3-6. 演習 3-A: ポートを数えて、失敗させる  (実測済み)",
        [
            (0, "① CPE に降りたプレフィックスから共有 IPv4 と PSID を確認し、3-3 の計算と照合", "b"),
            (0, "② ポート集合を紙に書き出し、飛び飛びであることを目で見る (240 個 = 16 × 15)", "b"),
            (0, "③ R3 (ポート制限) を再現する。OpenWrt が CE なら enforce は不要", "b"),
            (1, "for i in $(seq 20); do curl -s --limit-rate 1k -m 30 ... & done", "c"),
            (1, "紙に書いた 240 と違い、17 本目以降が失敗する (成功は 16 本)", "w"),
            (1, "nft の snat が終端判定で先頭ブロックのみ使用。実機では必ず数え直す", ""),
            (0, "④ R4: DS-Lite でポート開放を設定し、INET-SIM から叩いて届かないことを確認", "b"),
            (1, "前提: VNE と INET-SIM が別 VM (SPLIT_INET=1)。同居だと逆の結果になる", "w"),
            (0, "⑤ MTU の確認", "b"),
            (1, "ping -M do -s 1432 は通る / -s 1472 は通らない", "c"),
            (1, "5MB 転送の合否は run-checks.sh に任せる (curl 単体では遅いと落ちるの区別がつかない)", ""),
            (0, "⑥ 結果を test-matrix.md §5 実績表に記録する。記録しないと次の案件で活きない", "w"),
        ],
    )

    # ================= 第4部: 検証結果(実測) =================
    section_slide(prs, "第4部", "検証結果 — 本当にそうなるのか",
                  "ここから先は全部「測った値」です。推測と実測を必ず分けて書いてあります")

    bullet_slide(
        prs,
        "この部の読み方 — 「たぶん」と「測った」を分ける",
        [
            (0, "なぜこの部があるのか", "b"),
            (1, "設計どおりに動く、と思って客先に出た結果ハマるのを防ぐため", ""),
            (1, "自宅に検証ラボを組み、実機とシミュレータの両方で 1 つずつ確かめました", ""),
            (0, "何で測ったか(3 つの環境)", "b"),
            (1, "① 仮想ラボ (Proxmox 上に NGN / VNE / BRAS / INET を再現。CPE は OpenWrt)", ""),
            (1, "② 実機 (Cisco 892FJ を①のアクセス網に物理接続。クライアントは Windows 11 実機)", ""),
            (1, "③ CML (Cat8000v = 仮想ルータ。892FJ が MAP-E 非対応なので Cisco 実装の確認はこちら)", ""),
            (0, "覚えておいてほしい結論", "w"),
            (1, "トラブル再現レシピ 7 件のうち 5 件が「手順書どおりでは想定と違った」", ""),
            (1, "つまり、机上の想定は 7 割方どこかズレます。だから測る", ""),
        ],
    )

    diagram_slide(
        prs,
        "検証ラボの構成(これを会社の VMware にも持ち込みます)",
        [
            ("CPE", "OpenWrt\nまたは実機 892FJ"),
            ("NGN-SIM", "RA / DHCPv6-PD\nradvd + Kea"),
            ("VNE", "MAP-E BR\nDS-Lite AFTR"),
            ("INET-SIM", "Web + DNS\n模擬インターネット"),
        ],
        note=[
            "PPPoE 側 (BRAS) も別に立ててあるので、PPPoE → IPoE の切替そのものを演習できます。",
            "アドレスは全部ドキュメント用 (2001:db8::/32 / 198.51.100.0/24 / 203.0.113.0/24)。",
            "※ VNE と INET-SIM は必ず別 VM に分けること。同居させると DS-Lite の網側 NAT が効かず、",
            "  「疎通するので成功に見えるが、実は経路が違う」という一番タチの悪い状態になります。",
        ],
    )

    table_slide(
        prs,
        "実機 Cisco 892FJ で DS-Lite を通した結果(サイクル4)",
        ["確認したこと", "実測値", "何を示すか"],
        [
            ["出口 IPv4 アドレス", "203.0.113.1", "網側 (AFTR) で NAT された = DS-Lite が成立している"],
            ["ping 203.0.113.80", "0% 損失 / TTL=62", "Linux の初期値 64 から 2 減 = ルータを 2 つ越えている"],
            ["MTU 1460 (DF 付き 1432B)", "通過", "トンネル分を引いた MTU が確保できている"],
            ["MTU 1500 (DF 付き 1472B)", "192.168.100.1 から Frag Needed", "CPE 自身が生成 = PMTUD が正常に動いている"],
            ["5MB ダウンロード", "5,242,880 バイト / 1.15 秒", "MSS clamp が効いている (MTU ブラックホールなし)"],
            ["DNS", "A と AAAA を正しく応答", "ラボ DNS に到達している"],
        ],
        [3.6, 3.4, 5.0],
        footer="クライアントは実機 Windows 11。ip tcp adjust-mss 1420 を入れないと 5MB が完走しません",
        fsize=11,
    )

    bullet_slide(
        prs,
        "892FJ で分かった「手順書に書いておかないと必ず詰まる」3 点",
        [
            (0, "① RA 方式でないと autoconfig は何も取れない", "b"),
            (1, "ひかり電話ありの構成 (PD 方式) では radvd がプレフィックスを配らない", ""),
            (1, "症状: No global unicast address is configured。しかも RA 自体は届いている", ""),
            (1, "→ RA は見えるのにアドレスが付かない。原因に辿り着けない詰まり方をする", "w"),
            (0, "② classic IOS は MSS clamp が必須", "b"),
            (1, "ip tcp adjust-mss 1420 を入れないと、小さいページは見えるのに 5MB が止まる", ""),
            (1, "「Web は見えるがファイルが落ちない」の正体はほぼこれです", ""),
            (0, "③ CPE が WAN 側にも RA を撒く", "b"),
            (1, "ipv6 unicast-routing を入れると、CPE 自身がアクセス網に RA を流し始める", ""),
            (1, "実網でやると上流が CPE をルータとして学習する。ipv6 nd ra suppress all で止める", ""),
            (1, "LAN 側には入れないこと。配下のクライアントが SLAAC できなくなります", "w"),
        ],
        footer="詳細: docs/ipoe-lab/runbook-vmware.md §8 (実機手順) / §9-M・§9-N (詰まったとき)",
    )

    table_slide(
        prs,
        "MAP-E の計算は合っていたか(仮想ラボ・サイクル3)",
        ["項目", "計算で出した期待値", "実測", "判定"],
        [
            ["共有 IPv4 アドレス", "198.51.100.10", "198.51.100.10", "一致"],
            ["PSID", "5", "5", "一致"],
            ["使える先頭ポートブロック", "4176 - 4191", "4176 - 4191", "一致"],
            ["ポート総数 (設計値)", "15 ブロック × 16 = 240", "実効 16 (OpenWrt の癖)", "要注意"],
            ["MTU", "1460", "1460", "一致"],
        ],
        [3.4, 3.4, 3.0, 2.2],
        footer="計算どおりに出ます。ただしポート総数だけは CPE の実装に左右されます(次ページ)",
        fsize=12,
    )

    bullet_slide(
        prs,
        "Cisco 実装 (IOS XE) で 240 ポートを数えたかった — CML で試した結果",
        [
            (0, "なぜやったか", "b"),
            (1, "仮想ラボの OpenWrt では、設計 240 ポートに対し実効 16 しか使えなかった", ""),
            (1, "OpenWrt 固有の癖なのか、MAP-E 一般の話なのかが分からないままだった", ""),
            (1, "892FJ は classic IOS で MAP-E 非対応 → CML の Cat8000v (IOS XE 17.15.01a) で確認", ""),
            (0, "組んだ構成 (CML 上に 3 ノード)", "b"),
            (1, "client (Ubuntu) — Cat8000v (MAP-E CE) — BR 役 (Ubuntu)", ""),
            (1, "クライアントから送信元ポートを 400 個ばらけさせて UDP を撃ち、", ""),
            (1, "WAN 側リンクをキャプチャして「実際に何個の外側ポートが使われたか」を数える計画", ""),
            (0, "結果: 数えられませんでした。ただし別のことが 2 つ分かりました", "w"),
            (1, "① MAP-E の設定構文と、IOS XE が導出する値は設計と完全一致した", ""),
            (1, "② その設定を入れるとデータプレーンが落ちてルータが再起動する", ""),
        ],
    )

    table_slide(
        prs,
        "① IOS XE の MAP-E 設定 — 構文と導出値は設計どおりだった",
        ["項目", "こちらが与えた値", "IOS XE が導出した値", "設計値との一致"],
        [
            ["share-ratio", "256 (明示)", "Share-ratio 256 / bits 8", "一致 (PSID 長 8 ビット)"],
            ["start-port", "4096 (明示)", "Start-port 4096", "一致"],
            ["contiguous-ports", "与えていない", "Contiguous-ports 16 / bits 4", "一致 (自動導出)"],
            ["port-offset-bits", "与えていない", "Port-offset-bits 4", "一致 (自動導出)"],
            ["port-set-id", "5 (明示)", "Port-set-id 5", "一致 (PSID)"],
            ["ポート総数", "—", "15 ブロック × 16 = 240", "設計値 240 と一致"],
        ],
        [2.6, 2.8, 3.4, 3.2],
        footer="share-ratio と start-port だけ与えれば、残りは IOS XE が計算します。手計算と合いました",
        fsize=12,
    )

    bullet_slide(
        prs,
        "② 設定は通るのに、データプレーンが落ちる",
        [
            (0, "起きたこと", "b"),
            (1, "MAP-E の設定投入から約 60 秒後に転送エンジン (QFP) が落ちて再起動", ""),
            (1, "%PMAN-0-PROCFAILCRIT: critical process qfp_ucode_c8kv has failed (rc 134)", "c"),
            (0, "原因の切り分け (対照実験)", "b"),
            (1, "設定なし → 落ちない", ""),
            (1, "map-e domain はあるが port-parameters が構文エラーで入らなかった → 落ちない", ""),
            (1, "port-parameters が通って MAP-E が機能する状態になった → 落ちる", "w"),
            (1, "非対応と明記されている /64 の BMR を外し、正しい /56 の 1 本だけにしても落ちる", "w"),
            (0, "だから何が言えるか", "b"),
            (1, "CML の Cat8000v 17.15.01a では MAP-E の転送は確認できない (設定表示までは正しい)", ""),
            (1, "実機の IOS XE がどうかは、これでは分かりません。断定しないこと", ""),
            (0, "この件の一番の教訓", "w"),
            (1, "show コマンドが正しい値を返すことと、実際に転送できることは別物", ""),
            (1, "検証は必ず「パケットを流して測る」ところまでやる。設定確認で終わらせない", ""),
        ],
        footer="バックログ 10 (実効ポート数) は未解決のまま。実機 IOS XE か RTX 等の別 CE が要ります",
    )

    # ---------------- 壁 Q&A ----------------
    section_slide(prs, "つまずきの壁", "現場で実際に出る質問と答え方",
                  "現場で実際に出る質問と、その答え方",
                  notes=[
                      "【講師へ】",
                      "・**生徒役は若手にやらせること。**講師が両方やると聞き流される。",
                      "・時間割の外。各回の後ろに +15〜20 分で回す。",
                      "・全部やらない。その回で出た質問に近いものを 2〜3 個選ぶ。",
                  ])

    qa_slide(prs, 1, "IPoE って NTT がやってるサービス名ですか?",
             ["いや、IP over Ethernet の略。PPP を挟まずに Ethernet に直接 IP を載せる接続方式の名前。",
              "PPPoE は ID/パスワードでセッションを張る工程が入るが、IPoE はそれがない。",
              "つないだら RA か DHCPv6 でアドレスが降ってくる。"],
             ref="付録 A.3 (p.443) — PPPoE と IPoE が並べて書いてある")

    qa_slide(prs, 2, "「IPoE にすると速い」って営業が言うんですけど、なんで速いんですか?",
             ["IPoE 自体が速いわけじゃない。PPPoE は網終端装置 (NTE) を必ず通り、そこの収容数が決まっていて夕方混む。",
              "IPoE はそこを通らない。つまり「速い」の理由はプロトコルではなく通る箱が違うから。",
              "客への説明もこの言い方をする。"],
             dissent="それも常に正しくない。IPoE 側の VNE が混んでれば普通に遅い。「IPoE にしたのに遅い」案件は実在する。だから「速くなります」と言い切らないこと。約束すると後で詰められる。",
             ref="付録 A.3 (p.443)")

    qa_slide(prs, 3, "IPv6 にしたら IPv4 は使えなくなるんですか?",
             ["ならない。IPoE でも IPv4 は使える。ただし IPv6 のトンネルの中を通る (MAP-E / DS-Lite)。",
              "だから「IPv4 の通信が IPv6 に依存する」という新しい壊れ方が生まれる。",
              "IPv6 側が不調なら IPv4 も道連れになる、という点が PPPoE 時代と決定的に違う。"],
             ref="21.1 共存技術のバリエーション (p.381) / 第24章 (p.421)")

    qa_slide(prs, 4, "/64 とか /56 の数字は何ですか? IPv4 の /24 みたいなもの?",
             ["考え方は同じ (前から何ビットがネットワーク部か)。ただ運用の感覚が全然違う。",
              "IPv4 は「アドレスを 1 個もらう」、IPv6 は「プレフィックスを 1 個もらって中を自分で切る」。",
              "/56 をもらえば /64 が 256 個。SLAAC を使う限りホスト部は 64 ビット。"],
             ref="3.12 ユーザへの IPv6 アドレス割り当て (p.89) / 7.1 SLAAC の流れ (p.160)")

    qa_slide(prs, 5, "fe80:: のアドレスがいっぱい出てくるんですが、消していいですか?",
             ["消したら通信が全部死ぬ。リンクローカルアドレスといって、同じリンク内だけで使う。",
              "IPv6 のルーティングも近隣探索も、実はこれを使って動いている。RA の送信元も fe80::。",
              "IPv4 の 169.254.x.x とは役割が全然違う (あっちは DHCP 失敗時の敗戦処理)。"],
             ref="3.5 リンクローカルユニキャストアドレス (p.74) / 第6章 (p.127)")

    qa_slide(prs, 6, "tcpdump に ICMPv6 がめちゃくちゃ流れてるんですが、攻撃されてます?",
             ["正常。IPv6 では ICMPv6 がプロトコルの本体機能を担っている。",
              "ARP の代わり (NS/NA)、ルータ広告 (RS/RA)、Path MTU 通知、全部 ICMPv6。",
              "だから ICMPv6 を全部 drop するファイアウォール設定は IPv6 を壊す。"],
             dissent="逆に「全許可」も間違い。RFC 4890 相当の選別が必要。「ICMP を落とすと MTU 探索が壊れる」原理は v4/v6 共通で、このラボの R5 は IPv4 側 (ICMPv4 Fragmentation-Needed) で再現している。",
             ref="14.9 ICMPv6 を無条件にすべてフィルタリングすべきではない (p.279) / ラボ R5")

    qa_slide(prs, 7, "MAP-E の「ポートが制限される」って、何個使えるんですか?",
             ["方式のパラメータ次第。ラボの設定 (ea-len 16 / psid-len 8 / psid-offset 4) だと 240 個。",
              "16 ポート × 15 ブロックで、連続していない。計算は 3-2 / 3-3 のスライド。",
              "「4176〜8000 の範囲」のように連続だと思って NAT 設定を書くと動かない。"],
             dissent="240 は少ない。ブラウザは 1 タブで数十本張る。台数の多い拠点や監視系がいると足りなくなる。「MAP-E だから大丈夫」ではなく、拠点の同時セッション数を実測してから方式を決めること。",
             ref="24.3 A+P (p.423) / 24.4 (p.425) / ラボ R12 で必ず一度失敗させる")

    qa_slide(prs, 8, "客先で「ポート開放して」と言われました。IPoE でもできますよね?",
             ["方式を先に確認する。MAP-E なら自分のポート集合の中でのみ可 (80/443 は基本無理)。",
              "DS-Lite なら不可。NAT が網側 (AFTR) にあり、利用者からそれを操作する手段が提供されていない。",
              "代替は固定 IP 系サービスか PPPoE 併用。聞かずに「できます」と答えると案件が炎上する。"],
             ref="24.1 DS-Lite (p.421) / 19.4 IPv4 アドレス共有技術の課題 (p.362) / ラボ R4")

    qa_slide(prs, 9, "切替後に「Web は見えるけど添付ファイルが送れない」と言われました。",
             ["MTU を疑う。トンネルでヘッダが増えるので素の 1500 では通らない (引き算は 3-5)。",
              "小さいパケットは通るので「見えるけど大きいものが落ちる」という症状になる。",
              "一手は TCP MSS clamp か CPE の MTU 手動設定 (1460 等)。"],
             ref="第10章 (p.215) / 9.2 (p.208) / ラボ R5, R11 (/big.bin 5MB で再現)")

    qa_slide(prs, 10, "IPoE にしたらセキュリティは上がりますか?",
             ["下がる可能性がある。PPPoE + IPv4 NAT では「NAT があるから外から入れない」状態だったのが、",
              "IPoE で各端末にグローバル IPv6 が付くと、外から直接到達しうる。",
              "ファイアウォール設定が入っているかを必ず確認する。"],
             dissent="とはいえ「NAT はセキュリティ機能ではない」というのも教科書の立場。NAT に守られていたつもりだっただけ。この誤解をそのまま客に共有すると後々困る。",
             ref="14.1 IPv6 は IPv4 よりもセキュアというわけではない (p.263) / 14.7 (p.278) / ラボ R13")

    # ---------------- 参照系 ----------------
    table_slide(
        prs,
        "用語の対訳表 — 同じものが 3 通りの呼び方で出てくる",
        ["教科書の用語", "現場でよく言う", "機器 / 画面での表記"],
        [
            ["近隣探索プロトコル", "ND、ネイバーディスカバリ", "Neighbor Discovery / ipv6 nd"],
            ["Router Advertisement メッセージ", "RA、ルータ広告", "RA / ipv6 nd ra"],
            ["プレフィックス", "プレフィックス、接頭辞", "prefix"],
            ["SLAAC (IPv6 アドレスの自動設定)", "RA 方式、ステートレス", "SLAAC / autoconfig"],
            ["DHCPv6-PD", "PD 方式、プレフィックス委任", "IA_PD / ipv6 dhcp pool"],
            ["DUID", "ドゥーイド、クライアント識別子", "DUID-LL / DUID-LLT"],
            ["リンクローカルユニキャストアドレス", "リンクローカル、fe80", "link-local"],
            ["IPv4/IPv6 共存技術", "IPv4 over IPv6、v4v6", "transition technology"],
            ["MAP-E", "マップイー、v6 プラス系", "MAP-E / nat64 map-e (IOS XE)"],
            ["DS-Lite", "ディーエスライト、transix 系", "DS-Lite / AFTR"],
            ["A+P (ポートセット)", "ポート制限、PSID", "port-set / PSID"],
            ["CE / B4 / BR / AFTR", "宅内側 / 網側", "CE, B4 (宅内) / BR, AFTR (網側)"],
            ["Path MTU discovery", "PMTUD、MTU 探索", "path-mtu-discovery"],
        ],
        [4.3, 3.9, 4.0],
        footer="※ v6プラス / transix / クロスパス / IPv6オプション は「サービス名」。方式名ではない",
        fsize=11,
    )

    table_slide(
        prs,
        "トラブル再現レシピ → 読むべき節と、現場の対処",
        ["レシピ", "現象", "教科書", "対処"],
        [
            ["R1", "RA 方式なのに PD 前提設定 (最頻出)", "7.1-7.7 / 8.7", "ヒアリング#1 やり直し"],
            ["R2", "PPPoE セッション残留", "付録A.3", "ISP に強制切断依頼"],
            ["R3", "MAP-E ポート制限", "24.3 / 24.4 / 19.4", "使用ポートを範囲内に寄せる"],
            ["R4", "DS-Lite でポート開放不可", "24.1 / 19.3", "方式変更 / 固定IP / PPPoE併用"],
            ["R5", "MTU ブラックホール", "第10章 / 9.2 / 14.9", "MSS clamp / MTU 手動設定"],
            ["R6", "DNS フォールバック遅延", "17.2 / 15.3", "AAAA を止める / v6 経路を直す"],
            ["R7", "二重終端 (HGW + 配下ルータ)", "付録B / 13.2", "どちらで終端するか設計で確定"],
            ["R8", "v6 オプション未契約相当", "付録A.2", "契約状態確認"],
            ["R9", "プレフィックス変更で LAN が追従しない", "7.8 / 13.3", "固定前提の設計をやめる"],
            ["R10", "PPPoE 併用時の経路非対称", "13.1 / 13.2", "経路とアドレス選択を明示設定"],
            ["R11", "encaplimit 起因の片方向断", "4.3 / 22.7 / 14.10", "encaplimit none にする"],
            ["R12", "ポートセット非連続を無視した NAT 設定", "24.3 / 24.4", "ポート集合を計算し直す"],
            ["R13", "IPoE 化による外→内 IPv6 直接着信", "第14章 / 14.7", "CPE の v6 FW を確認・有効化"],
        ],
        [1.0, 5.0, 2.9, 3.3],
        footer="レシピ本体: docs/ipoe-lab/test-matrix.md §4 (全 13 種)",
        fsize=11,
    )

    bullet_slide(
        prs,
        "教科書とラボの関係 (MAP-E のインタフェース識別子)",
        [
            (0, "第2版第2刷の更新履歴に「MAP-E のインタフェース識別子を RFC 7597 準拠に」修正の記載がある", ""),
            (0, "ラボの既定も RFC 7597。つまり教科書と一致している", "b"),
            (1, "setup-map-br.sh の CE_MAP_ADDR は RFC 7597 の IID", ""),
            (0, "食い違うのは実サービス側", "b"),
            (1, "日本の VNE は draft-ietf-softwire-map-03 互換", ""),
            (1, "実案件のリハーサル時だけ option legacymap '1' を立て、CE_MAP_ADDR を再計算する", ""),
            (0, "ラボ既定のまま legacymap '1' を入れると BR とアドレスが合わず MAP-E が全断する", "w"),
            (0, "伝え方: 「教科書どおり (RFC 7597) で動くのがラボ。実サービスは draft-03」", "b"),
            (1, "RFC = 確定した仕様、draft = 確定前の草案", ""),
            (1, "日本の VNE は草案のまま製品化されて普及した、という経緯", ""),
            (0, "詳細: build.md §3 のコメント / research-notes.md §2", ""),
        ],
        footer="こういう食い違いは IPv6 の世界では珍しくない。だから検証環境が必要になる",
    )

    bullet_slide(
        prs,
        "修了確認 — 理由付きで答えられるか (5問)",
        [
            (0, "1. ひかり電話ありの案件。CPE の IPv6 設定はどうするか。その前に確認すべきことは何か", ""),
            (1, "(HGW の有無と型番 → HGW と配下ルータのどちらで終端するかを設計で決める。", ""),
            (1, " ひかり電話あり → 即 PD 設定は二重終端 R7 を作る)", ""),
            (0, "2. 「Minecraft サーバ (TCP 25565) を公開したい」。まず何を確認するか", ""),
            (1, "MAP-E だったとき、25565 が自分の 240 ポートの集合に入る見込みはどれくらいか", ""),
            (0, "3. 「切替後、Web は見えるが大きいファイルが送れない」。何を疑い、どう直すか", ""),
            (0, "4. ICMPv6 を全部 deny した。何が壊れるか 2 つ挙げよ", ""),
            (0, "5. IPoE 化で PC に付いたグローバル IPv6。セキュリティ上、何を確認するか", ""),
            (0, "答えの暗記が目的ではない。「どの節・どのファイルを見ればいいか」が分かれば合格", "b"),
            (1, "現場では調べる時間がある。調べる先が分からないことが事故を生む", ""),
        ],
    )

    bullet_slide(
        prs,
        "参考資料 / 出典",
        [
            (0, "教科書", "b"),
            (1, "小川晃通『プロフェッショナル IPv6 第2版』ラムダノート, 2021年12月", ""),
            (1, "(第2刷 2023年3月) CC BY-NC-SA", ""),
            (1, "※ 節番号は刷でずれないが、ページ番号はずれ得る。参照は節番号+節タイトルを主キーに", "w"),
            (1, "※ 教科書の本文・図版はこのスライドに貼らないこと (NC / SA 条項)。社内限定資料", "w"),
            (0, "社内資料 (本リポジトリ)", "b"),
            (1, "docs/ipoe-lab/study-guide.md … 本スライドの正本。対応表と計算例の全量はこちら", ""),
            (1, "docs/ipoe-lab/README.md … ラボ全体設計・アドレス計画 (パラメータの正本)", ""),
            (1, "docs/ipoe-lab/build.md … 構築手順・モード別期待値表 (パラメータの正本)", ""),
            (1, "docs/ipoe-lab/test-matrix.md … ヒアリングシート / マトリクス / レシピ R1-R13", ""),
            (1, "docs/ipoe-lab/research-notes.md … 実サービスのパラメータ、既知の落とし穴", ""),
            (1, "docs/ipoe-lab/runbook-vmware.md … 社内 VMware 用 (執筆前)", ""),
            (0, "このスライドの再生成 (リポジトリのルートから)", "b"),
            (1, "python3 docs/ipoe-lab/slides/build-toranomaki.py", "c"),
        ],
    )

    prs.save(OUT)
    print(f"saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
